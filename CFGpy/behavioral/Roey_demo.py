import os

import pandas as pd
from CFGpy.behavioral import Downloader, Parser, Preprocessor
from CFGpy.behavioral.MRIParser import MRIParser
from CFGpy.behavioral.Preprocessor import DEFAULT_OUTPUT_FILENAME
from CFGpy.behavioral.data_structs import PreprocessedPlayerData, PATH_FROM_REP_ROOT
from CFGpy.behavioral._consts import *
from pptx import Presentation
from pptx.util import Inches
import json # Roey added for saving the vanilla after more preprocessing

CSV_URL = "https://api.creativeforagingtask.com/v1/event.csv?game=4cb46367-7555-\
42cb-8915-152c3f3efdfb&entityType=event&after=2021-05-23T10:51:00.\
000Z"  # link Roey sent
CSV_FILE_PATH = "/home/roey/Documents/CFG_data/event.csv"
ROY_TEST_JASON = "/home/roey/PycharmProjects/CFGpy/CFGpy/behavioral/test_file1.json"

PPT_OUTPUT_PATH = 'output/'

def __from_raw_data(raw_data):
    parser = MRIParser(raw_data)
    print("Parsing...")
    parsed_data = parser.parse()

    preprocessor = Preprocessor(parsed_data)
    print("Segmenting...")

    preprocessor.dump()  # save JSON

    return preprocessor.preprocess()


def from_url(red_metrics_csv_url):
    downloader = Downloader(red_metrics_csv_url)
    print("Downloading raw data...")
    raw_data = downloader.download()
    return __from_raw_data(raw_data)


def from_file(red_metrics_csv_path):
    raw_data = pd.read_csv(red_metrics_csv_path)
    return __from_raw_data(raw_data)


def from_json(jason_path):
    preprocessor = Preprocessor.from_json(jason_path)
    return preprocessor.preprocess()


def add_one_plot_to_ppt(prs, image_path, title=None):
    # Add a slide with a title and content layout
    slide_layout = prs.slide_layouts[5]  # Choosing a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Define image placement on slide
    left = Inches(1)
    top = Inches(1)
    height = Inches(4.5)

    # Add image to slide
    slide.shapes.add_picture(image_path, left, top, height=height)

    # Add a title to the slide
    if title is not None:
        title_box = slide.shapes.title
        title_box.text = title

def add_both_plots_to_ppt(prs, shapesImagePath, plotGalleryImagePath, title):
    slide_layout = prs.slide_layouts[5]  # Choosing a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add gallery plot image
    left = Inches(0.1)
    top = Inches(2)
    height = Inches(3)
    slide.shapes.add_picture(plotGalleryImagePath, left, top, height=height)

    # Add shapes image
    left = Inches(6)
    slide.shapes.add_picture(shapesImagePath, left, top, height=height)


def add_error_slide_to_ppt(prs, errors):
    """
    Adds a slide with a list of error messages to the presentation.
    :param prs: Presentation object.
    :param errors: List of error messages.
    """
    slide_layout = prs.slide_layouts[1]  # Choosing a title and content slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.title
    title_box.text = "Errors"

    # Add error messages
    content_box = slide.placeholders[1]
    content_box.text = "\n".join(errors)


def create_players_cluster_times_csv(preprocessed_data, output_folder=PATH_FROM_REP_ROOT+"clusters_times_csvs/"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for player_data in preprocessed_data:
        if not player_data[PARSED_PLAYER_ID_KEY].startswith("9999"):  # skip the non-player data
            data = PreprocessedPlayerData(player_data)

            # explore_times, exploit_times = data.get_cluster_times()
            #
            player_id = player_data[PARSED_PLAYER_ID_KEY]
            csv_file_path = os.path.join(output_folder, f"Player_{player_id}_cluster_times.csv")

            # Save DataFrame to CSV
            data.clusters_times_to_csv(csv_file_path)
            #
            # def create_data_frame(phase, cluster_times, columns=("start", "end")):
            #     start, end = columns[0], columns[1]
            #     df = pd.DataFrame(cluster_times, columns=[start,end])
            #     df["phase"] = phase
            #     return df
            #
            # explore_df = create_data_frame(EXPLORE_KEY, explore_times)
            # exploit_df = create_data_frame(EXPLOIT_KEY, exploit_times)
            #
            # # # Create DataFrame for explore times
            # # explore_df = pd.DataFrame(explore_times, columns=["gallery_out_time", "gallery_in_time"])
            # # explore_df["phase"] = EXPLORE_KEY
            # #
            # # # Create DataFrame for exploit times
            # # exploit_df = pd.DataFrame(exploit_times, columns=["gallery_out_time", "gallery_in_time"])
            # # exploit_df["phase"] = EXPLOIT_KEY
            #
            # # Combine both DataFrames
            # cluster_times_df = pd.concat([explore_df, exploit_df], ignore_index=True)
            #
            # # Save DataFrame to CSV
            # cluster_times_df.to_csv(csv_file_path, index=False)

def create_empty_moves_csv(preprocessed_data, output_folder=PATH_FROM_REP_ROOT+"empty_moves_csvs/"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for player_data in preprocessed_data:
        if not player_data[PARSED_PLAYER_ID_KEY].startswith("9999"):  # skip the non-player data
            data = PreprocessedPlayerData(player_data)

            player_id = player_data[PARSED_PLAYER_ID_KEY]
            csv_file_path = os.path.join(output_folder, f"Player_{player_id}_cluster_times.csv")

            # Save DataFrame to CSV
            data.get_empty_moves_times(csv_file_path)

def create_players_game_presentation(preprocessed_data, delta_t=True):
    # Create a presentation object
    prs = Presentation()
    # Collect error messages
    error_messages = []
    # Plot each subject's game OR add to presentation
    for player_data in preprocessed_data:

        if not player_data[PARSED_PLAYER_ID_KEY].startswith("9999"):  # skip the non-player data
            data = PreprocessedPlayerData(player_data)
            shapesImagePath = data.plot_clusters()

            if delta_t:
                plotGalleryImagePath = data.plot_gallery_dt()
            else:
                plotGalleryImagePath = data.plot_gallery_steps_over_dt() #TODO: Consider adding a flag for plot_gallery_steps

            if shapesImagePath != -1 and plotGalleryImagePath != -1:
                add_both_plots_to_ppt(prs, shapesImagePath, plotGalleryImagePath,
                                      f"Player {player_data[PARSED_PLAYER_ID_KEY]}")
                os.remove(shapesImagePath)  # Clean up the image file
                os.remove(plotGalleryImagePath)  # Clean up the image file
            else:
                error_messages.append(
                    f"{player_data[PARSED_PLAYER_ID_KEY]}, missing necessary phases --> can't create plot")

    # Add a single error slide if there are any errors
    if error_messages:
        add_error_slide_to_ppt(prs, error_messages)

    return prs

def create_all_shapes_presentation(preprocessed_data):
    # Create a presentation object
    prs = Presentation()
    # Collect error messages
    error_messages = []
    # Plot each subject's game OR add to presentation
    for player_data in preprocessed_data:

        if not player_data[PARSED_PLAYER_ID_KEY].startswith("9999"):  # skip the non-player data
            data = PreprocessedPlayerData(player_data)

            shapesImagePath = data.plot_shapes()

            if shapesImagePath != -1:
                add_one_plot_to_ppt(prs, shapesImagePath, f"Player {player_data[PARSED_PLAYER_ID_KEY]}")
                os.remove(shapesImagePath)  # Clean up the image file
            else:
                error_messages.append(
                    f"{player_data[PARSED_PLAYER_ID_KEY]}, missing necessary phases --> can't create plot")

    # Add a single error slide if there are any errors
    if error_messages:
        add_error_slide_to_ppt(prs, error_messages)

    return prs

if __name__ == '__main__':
    # NOTE: Change here to one of 3 options for loading the data

    # Load vanilla to preprocess it with the MRI algorithm and dump it into a new json file
    #preprocessed_data = from_json("/Volumes/HartLabNAS/Projects/CFG/vanilla_data/vanilla.json")
    #path = "/Volumes/HartLabNAS/Projects/CFG/vanilla_data/vanilla_shuffled_noMRIFix.json"
    #with open(path, "w") as out_file:
    #    json.dump(preprocessed_data, out_file)

    # Load data from url & save new JSON
    #preprocessed_data = from_url(CSV_URL)

    # Load data from csv
    #preprocessed_data = from_file(CSV_FILE_PATH)

    # Load processed data from json
    preprocessed_data = from_json(DEFAULT_OUTPUT_FILENAME)

    # Sort by ID
    preprocessed_data = sorted(preprocessed_data, key=lambda x: x["id"])  # sort by subjects' ID

    def create_presentation_with_all_shapes_plot():
        prs_all_shapes =create_all_shapes_presentation(preprocessed_data)
        prs_all_shapes_name = f"{PPT_OUTPUT_PATH}all_shapes_presentation.pptx"
        prs_all_shapes.save(prs_all_shapes_name)
        print(f"PowerPoint presentation saved as '{prs_all_shapes_name}'")

    def create_presentation_with_both_plots():
        plot_by_delta_t = True # change this from False --> True: for choosing weather it'll create presentation by delta_t or steps
        prs_game = create_players_game_presentation(preprocessed_data, delta_t=plot_by_delta_t)
        if not plot_by_delta_t:
            prs_game_name = f"{PPT_OUTPUT_PATH}games_presentation_efficiency{MIN_EFFICIENCY_FOR_EXPLOIT}_paceSpecific_5mad_rmvEmptyTimeAllSteps_efficiencyLessThan1DELETE.pptx"
        else:
            prs_game_name = f"{PPT_OUTPUT_PATH}games_presentation_efficiency{MIN_EFFICIENCY_FOR_EXPLOIT}_paceSpecific_5mad_rmvEmptyTimeAllSteps_efficiencyLessThan1_delta_tDELETE.pptx"

        prs_game.save(prs_game_name)
        print(f"PowerPoint presentation saved as '{prs_game_name}'")

    def create_players_cluster_times():
        output_folder = PATH_FROM_REP_ROOT+"clusters_times_csvs/"
        create_players_cluster_times_csv(preprocessed_data, output_folder)
        print(f"csv's with players cluster saved in folder: '{output_folder}'")

    def create_players_empty_moves_times():
        output_folder = PATH_FROM_REP_ROOT+"empty_moves_csvs/"
        create_empty_moves_csv(preprocessed_data, output_folder)
        print(f"csv's with players cluster saved in folder: '{output_folder}'")

    # create_presentation_with_both_plots()
    create_players_cluster_times()
    create_players_empty_moves_times()